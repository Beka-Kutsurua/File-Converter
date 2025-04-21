from docx2pdf import convert
from pdf2docx import Converter
import os
import fitz
import pythoncom
import shutil
from PIL import Image
import tempfile
from pptx import Presentation
from pptx.util import Inches

def copy_file(input_path, output_path):
    shutil.copyfile(input_path, output_path)

def docx_to_pdf(input_path, output_path):
    pythoncom.CoInitialize()
    try:
        convert(input_path, output_path)
    finally:
        pythoncom.CoUninitialize()
    return output_path

def pdf_to_docx(input_path, output_path):
    converter = Converter(input_path)
    converter.convert(output_path, start=0, end=None)
    converter.close()

def pdf_to_jpg(input_pdf_path, output_dir):
    doc = fitz.open(input_pdf_path)
    output_files = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=200)
        output_path = os.path.join(output_dir, f"page_{page_num + 1}.jpg")
        pix.save(output_path)
        output_files.append(output_path)

    doc.close()
    return output_files

def pdf_to_png(input_pdf_path, output_dir):
    doc = fitz.open(input_pdf_path)
    output_files = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=200)
        output_path = os.path.join(output_dir, f"page_{page_num + 1}.png")
        pix.save(output_path)
        output_files.append(output_path)

    doc.close()
    return output_files


def images_to_pdf(input_paths, output_path):
    images = []
    for path in input_paths:
        img = Image.open(path)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        images.append(img)
    if images:
        images[0].save(output_path, save_all=True, append_images=images[1:], format='PDF')

def jpg_to_png(input_paths, output_dir):
    output_files = []
    for path in input_paths:
        img = Image.open(path)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        
        base = os.path.splitext(os.path.basename(path))[0]
        filename = base + ".png"
        output_path = os.path.join(output_dir, filename)

        print(f"[jpg_to_png] Saving to: {output_path}")
        img.save(output_path, format="PNG")
        output_files.append(output_path)

    return output_files

def png_to_jpg(input_paths, output_dir):
    output_files = []
    for path in input_paths:
        img = Image.open(path)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        
        base = os.path.splitext(os.path.basename(path))[0]
        filename = base + ".jpg"
        output_path = os.path.join(output_dir, filename)

        print(f"[png_to_jpg] Saving to: {output_path}")
        img.save(output_path, format="JPEG")
        output_files.append(output_path)

    return output_files

def pdf_to_pptx(input_pdf_path, output_path):
    doc = fitz.open(input_pdf_path)
    presentation = Presentation()
    blank_slide_layout = presentation.slide_layouts[6]

    temp_images = []  # Track temp files to delete later

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(dpi=150)

        tmp_img_path = os.path.join(tempfile.gettempdir(), f"pdf_slide_{page_num}.png")
        pix.save(tmp_img_path)
        temp_images.append(tmp_img_path)

        slide = presentation.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(tmp_img_path, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))

    doc.close()
    presentation.save(output_path)

    for path in temp_images:
        try:
            os.remove(path)
        except Exception as e:
            print(f"Warning: Could not delete temp file {path}: {e}")
    
    return output_path


